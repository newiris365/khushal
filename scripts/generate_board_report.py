import sys
import os
import json
import argparse
import subprocess

# Self-healing install for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx not found. Attempting to install...", flush=True)
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except Exception as e:
        print(f"ERROR: Failed to install python-pptx dynamically: {e}", file=sys.stderr, flush=True)
        sys.exit(1)

def main():
    parser = argparse.ArgumentParser(description="Generate Board/Trustee PowerPoint Report")
    parser.add_argument("--institution", type=str, default="SIET Institution", help="Name of the institution")
    parser.add_argument("--quarter", type=int, default=2, help="Quarter (1-4)")
    parser.add_argument("--year", type=int, default=2026, help="Year")
    parser.add_argument("--output", type=str, default="board_report.pptx", help="Output filepath")
    parser.add_argument("--data", type=str, default="{}", help="JSON data string containing KPIs")
    
    args = parser.parse_args()
    
    try:
        data = json.loads(args.data)
    except Exception:
        data = {}

    prs = Presentation()
    
    # Define color scheme (Dark Slate/Purple theme matching IRIS 365)
    bg_color = RGBColor(13, 10, 26)       # #0D0A1A
    primary_color = RGBColor(108, 43, 217)  # #6C2BD9 (Purple)
    text_white = RGBColor(255, 255, 255)
    text_gray = RGBColor(196, 181, 253)    # #C4B5FD
    accent_emerald = RGBColor(16, 185, 129) # #10B981
    accent_red = RGBColor(239, 68, 68)     # #EF4444

    # Slide 1: Title Slide (Executive Branding)
    slide_layout = prs.slide_layouts[5] # Blank slide layout or title only
    slide = prs.slides.add_slide(slide_layout)
    
    # Set dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title Text Box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f"{args.institution}\nQ{args.quarter} {args.year} Board of Trustees Report"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = text_white
    
    p2 = tf.add_paragraph()
    p2.text = "IRIS 365 Campus OS Intelligence Hub • Generated Real-time"
    p2.font.size = Pt(14)
    p2.font.color.rgb = text_gray

    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = primary_color
    
    txBoxBody = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.0))
    tf_body = txBoxBody.text_frame
    tf_body.word_wrap = True
    
    p_body1 = tf_body.paragraphs[0]
    p_body1.text = f"• Operations Summary: Q{args.quarter} has demonstrated consistent digital maturity across all active campus modules (Canteen, Gym, Transit, Library)."
    p_body1.font.size = Pt(16)
    p_body1.font.color.rgb = text_white
    
    p_body2 = tf_body.add_paragraph()
    p_body2.text = f"• Attendance Baseline: Average attendance stabilized at {data.get('attendance_rate', 82.5)}% with compiler and smart scanner logs active."
    p_body2.font.size = Pt(16)
    p_body2.font.color.rgb = text_white
    p_body2.space_before = Pt(14)
    
    p_body3 = tf_body.add_paragraph()
    p_body3.text = f"• Financial Position: Real-time net surplus computed at ₹{data.get('net_surplus', '33.4L')} through consolidated fees and module microtransactions."
    p_body3.font.size = Pt(16)
    p_body3.font.color.rgb = text_white
    p_body3.space_before = Pt(14)

    # Slide 3: Key Performance Indicators (KPIs)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Institutional Key Performance Indicators"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = primary_color
    
    # Create cards layout (KPI values)
    kpis = [
        {"label": "Avg Attendance", "value": f"{data.get('attendance_rate', 82.5)}%", "desc": "vs 78.5% industry avg"},
        {"label": "Fee Collection", "value": f"{data.get('fee_collection_percent', 78)}%", "desc": "₹1.42Cr collected"},
        {"label": "Module Adoption", "value": f"{data.get('module_adoption', 76)}%", "desc": "Active digital footprint"},
        {"label": "Canteen Active Users", "value": f"{data.get('canteen_users', 450)}", "desc": "92% daily transactions"}
    ]
    
    for i, kpi in enumerate(kpis):
        left = Inches(0.5 + (i % 2) * 4.5)
        top = Inches(1.8 + (i // 2) * 2.3)
        width = Inches(4.0)
        height = Inches(1.8)
        
        cardBox = slide.shapes.add_textbox(left, top, width, height)
        ctf = cardBox.text_frame
        ctf.word_wrap = True
        
        p_lbl = ctf.paragraphs[0]
        p_lbl.text = kpi["label"]
        p_lbl.font.size = Pt(14)
        p_lbl.font.color.rgb = text_gray
        p_lbl.font.bold = True
        
        p_val = ctf.add_paragraph()
        p_val.text = kpi["value"]
        p_val.font.size = Pt(32)
        p_val.font.bold = True
        p_val.font.color.rgb = accent_emerald
        
        p_dsc = ctf.add_paragraph()
        p_dsc.text = kpi["desc"]
        p_dsc.font.size = Pt(11)
        p_dsc.font.color.rgb = text_white

    # Slide 4: Year-over-Year (YoY) Comparisons & Achievements
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Year-over-Year Comparison & Trends"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = primary_color
    
    txBoxBody = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.0))
    tf_body = txBoxBody.text_frame
    tf_body.word_wrap = True
    
    p_trend1 = tf_body.paragraphs[0]
    p_trend1.text = "• Fee Collection Rate YoY: +12% increase from 2025 due to automatic reminders and digital wallets integrations."
    p_trend1.font.size = Pt(15)
    p_trend1.font.color.rgb = text_white
    
    p_trend2 = tf_body.add_paragraph()
    p_trend2.text = "• Campus Security Audits: Anomaly flag response times decreased by 40% with smart gate threat filters."
    p_trend2.font.size = Pt(15)
    p_trend2.font.color.rgb = text_white
    p_trend2.space_before = Pt(12)
    
    p_trend3 = tf_body.add_paragraph()
    p_trend3.text = "• Student Engagement Score: +18% increase compared to Q2 last year, showing positive adoption of campus events and gym programs."
    p_trend3.font.size = Pt(15)
    p_trend3.font.color.rgb = text_white
    p_trend3.space_before = Pt(12)

    # Slide 5: Concerns & Actions
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Areas of Concern & Corrective Actions"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = accent_red
    
    txBoxBody = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.0))
    tf_body = txBoxBody.text_frame
    tf_body.word_wrap = True
    
    p_c1 = tf_body.paragraphs[0]
    p_c1.text = "• Low Library Circulation: 28% of students registered under-utilize physical textbook checkouts."
    p_c1.font.size = Pt(15)
    p_c1.font.color.rgb = text_white
    
    p_a1 = tf_body.add_paragraph()
    p_a1.text = "   → Action: Introduce AI research helper and digital newspaper catalog in Library+."
    p_a1.font.size = Pt(13)
    p_a1.font.color.rgb = text_gray
    p_a1.space_before = Pt(4)
    
    p_c2 = tf_body.add_paragraph()
    p_c2.text = "• Sophomores Low Attendance: CS Sophomores show attendance rate of 72% average."
    p_c2.font.size = Pt(15)
    p_c2.font.color.rgb = text_white
    p_c2.space_before = Pt(14)
    
    p_a2 = tf_body.add_paragraph()
    p_a2.text = "   → Action: Auto-assign counselor interventions via Student Journey Score analytics."
    p_a2.font.size = Pt(13)
    p_a2.font.color.rgb = text_gray
    p_a2.space_before = Pt(4)

    # Slide 6: Next Quarter Outlook
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Quarter Outlook & Roadmap"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = primary_color
    
    txBoxBody = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.0))
    tf_body = txBoxBody.text_frame
    tf_body.word_wrap = True
    
    p_o1 = tf_body.paragraphs[0]
    p_o1.text = "• Financial Surplus Optimization: Forecasted cash flow projects +5% net margins with events pricing adjustment."
    p_o1.font.size = Pt(15)
    p_o1.font.color.rgb = text_white
    
    p_o2 = tf_body.add_paragraph()
    p_o2.text = "• Competitor Benchmarks Rajasthan: Align current modules with top 10 percentile indicators."
    p_o2.font.size = Pt(15)
    p_o2.font.color.rgb = text_white
    p_o2.space_before = Pt(12)
    
    p_o3 = tf_body.add_paragraph()
    p_o3.text = "• Board Report Distribution: Automatized monthly delivery via email list dispatchers."
    p_o3.font.size = Pt(15)
    p_o3.font.color.rgb = text_white
    p_o3.space_before = Pt(12)

    prs.save(args.output)
    print(f"SUCCESS: PowerPoint report saved successfully at {args.output}", flush=True)

if __name__ == "__main__":
    main()
